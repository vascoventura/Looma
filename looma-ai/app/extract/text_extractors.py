from pathlib import Path
import json
import re
import subprocess
import tempfile

import pymupdf
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup


def _looks_like_scanned_or_empty(text: str) -> bool:
    if not text:
        return True

    stripped = ''.join(ch for ch in text if not ch.isspace())
    if len(stripped) < 40:
        return True

    letters = sum(ch.isalpha() for ch in stripped)
    ratio = letters / max(len(stripped), 1)
    return ratio < 0.15


def _tesseract_ocr_image(image_path: Path, *, langs: str) -> str:
    with tempfile.TemporaryDirectory() as tmpdir:
        outbase = str(Path(tmpdir) / 'ocr')
        subprocess.run(
            ['tesseract', str(image_path), outbase, '-l', langs],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        txt_path = Path(outbase + '.txt')
        return txt_path.read_text(encoding='utf-8', errors='ignore') if txt_path.exists() else ''


def extract_pdf(
    path: Path,
    *,
    ocr_langs: str = 'eng+nep',
    ocr_if_low_text: bool = True,
    dpi: int = 200,
    max_ocr_pages: int | None = None,
):
    doc = pymupdf.open(path)
    pages = []

    ocr_pages = 0
    with tempfile.TemporaryDirectory() as tmpdir:
        tmpdir_path = Path(tmpdir)

        for i, page in enumerate(doc, start=1):
            text = page.get_text('text') or ''

            should_ocr = ocr_if_low_text and _looks_like_scanned_or_empty(text)
            if should_ocr and (max_ocr_pages is None or ocr_pages < max_ocr_pages):
                try:
                    pix = page.get_pixmap(dpi=dpi)
                    img_path = tmpdir_path / f'page_{i}.png'
                    pix.save(str(img_path))
                    ocr_text = _tesseract_ocr_image(img_path, langs=ocr_langs)
                    if ocr_text and not _looks_like_scanned_or_empty(ocr_text):
                        text = ocr_text
                        ocr_pages += 1
                except FileNotFoundError:
                    # tesseract not installed
                    pass
                except Exception:
                    pass

            pages.append({'page': i, 'text': text})

    return pages


def extract_docx(path: Path):
    doc = Document(path)
    text = '\n'.join(p.text for p in doc.paragraphs if p.text and p.text.strip())
    return [{'page': 1, 'text': text}]


def extract_pptx(path: Path):
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                parts.append(shape.text)
        slides.append({'page': i, 'text': '\n'.join(parts)})
    return slides


def extract_txt(path: Path):
    return [{'page': 1, 'text': path.read_text(encoding='utf-8', errors='ignore')}]


def extract_html(path: Path):
    raw = path.read_text(encoding='utf-8', errors='ignore')
    soup = BeautifulSoup(raw, 'html.parser')
    text = soup.get_text('\n')
    return [{'page': 1, 'text': text}]


def extract_vtt(path: Path):
    raw = path.read_text(encoding='utf-8', errors='ignore')
    # Strip WEBVTT header + timestamps/cue settings.
    raw = re.sub(r'^\ufeff?WEBVTT.*?$', ' ', raw, flags=re.M)
    raw = re.sub(r'^\d+$', ' ', raw, flags=re.M)
    raw = re.sub(r'^\d{2}:\d{2}:\d{2}\.\d{3}\s+-->\s+\d{2}:\d{2}:\d{2}\.\d{3}.*?$', ' ', raw, flags=re.M)
    raw = re.sub(r'^\d{2}:\d{2}\.\d{3}\s+-->\s+\d{2}:\d{2}\.\d{3}.*?$', ' ', raw, flags=re.M)
    text = re.sub(r'\s+', ' ', raw).strip()
    return [{'page': 1, 'text': text}]


def extract_image_ocr(path: Path, langs='eng+nep'):
    try:
        text = _tesseract_ocr_image(path, langs=langs)
    except FileNotFoundError:
        text = ''
    return [{'page': 1, 'text': text}]


def extract_keywords(path: Path):
    """`.keywords` files are JSON arrays of vocabulary entries used by Looma:
        [{"en": "father", "np": "बाबा", "def": "A male parent."}, ...]
    Flatten them to a readable list so the embeddings learn each vocabulary
    item alongside its translation and definition.
    """
    raw = path.read_text(encoding='utf-8', errors='ignore')
    try:
        data = json.loads(raw)
    except Exception:
        # Some legacy files have stray characters / concatenated JSON; fall back
        # to treating them as plain text so we still index whatever is there.
        return [{'page': 1, 'text': raw}]

    if not isinstance(data, list):
        return [{'page': 1, 'text': raw}]

    lines = []
    for entry in data:
        if not isinstance(entry, dict):
            continue
        en   = str(entry.get('en')   or '').strip()
        np_  = str(entry.get('np')   or '').strip()
        defn = str(entry.get('def')  or entry.get('definition') or '').strip()
        if not (en or np_ or defn):
            continue
        if en and np_:
            head = f'{en} — {np_}'
        else:
            head = en or np_
        lines.append(f'{head}: {defn}' if defn else head)
    return [{'page': 1, 'text': '\n'.join(lines)}]


def extract_geojson(path: Path):
    """`.geojson` files store map features with rich text properties (name,
    country, capital, currency, language, area, GDP, …). Flatten each feature's
    properties into 'key: value' lines so place names and country facts become
    semantically searchable.
    """
    raw = path.read_text(encoding='utf-8', errors='ignore')
    try:
        data = json.loads(raw)
    except Exception:
        return [{'page': 1, 'text': raw}]

    if isinstance(data, dict) and data.get('type') == 'FeatureCollection':
        features = data.get('features') or []
    elif isinstance(data, dict) and data.get('type') == 'Feature':
        features = [data]
    elif isinstance(data, list):
        features = data
    else:
        features = []

    lines = []
    for feat in features:
        if not isinstance(feat, dict):
            continue
        props = feat.get('properties') if isinstance(feat.get('properties'), dict) else {}
        if not props:
            continue
        parts = []
        for k, v in props.items():
            if v is None or v == '' or isinstance(v, (list, dict)):
                continue
            if isinstance(v, bool):
                parts.append(f'{k}: {str(v).lower()}')
            elif isinstance(v, (int, float)):
                parts.append(f'{k}: {v}')
            else:
                s = str(v).strip()
                if s:
                    parts.append(f'{k}: {s}')
        if parts:
            lines.append('; '.join(parts))
    return [{'page': 1, 'text': '\n'.join(lines) if lines else raw}]


def extract_any(path: Path, ocr_langs='eng+nep'):
    suffix = path.suffix.lower()

    if suffix == '.pdf':
        return extract_pdf(path, ocr_langs=ocr_langs)
    if suffix == '.docx':
        return extract_docx(path)
    if suffix == '.pptx':
        return extract_pptx(path)
    # Plain-text / Markdown-like Looma curriculum metadata. .summary / .outline /
    # .quiz / .objectives / .plan / .lesson live next to chapter PDFs in
    # content/chapters and carry chapter summaries, outlines, quizzes and lesson
    # plans — high-signal content for semantic search.
    if suffix in {'.txt', '.md', '.summary', '.outline', '.quiz',
                  '.objectives', '.plan', '.lesson'}:
        return extract_txt(path)
    if suffix == '.keywords':
        return extract_keywords(path)
    if suffix == '.geojson':
        return extract_geojson(path)
    if suffix == '.vtt':
        return extract_vtt(path)
    if suffix in {'.html', '.htm'}:
        return extract_html(path)
    if suffix in {'.png', '.jpg', '.jpeg', '.tif', '.tiff', '.bmp', '.webp'}:
        return extract_image_ocr(path, langs=ocr_langs)

    raise ValueError(f'Unsupported file type: {path.suffix}')
